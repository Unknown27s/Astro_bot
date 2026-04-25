"""
IMS AstroBot — Document Parsers
Extracts text from PDF, DOCX, TXT, XLSX, CSV, PPTX, and HTML files.

PDF improvements:
  - pdfplumber with explicit table strategies (lines + text-based)
  - Table bounding boxes used to exclude table regions from prose text
    (prevents the same content appearing twice)
  - Header row detection and markdown-style formatting
  - Merged / empty cell handling
  - Three-strategy fallback: pdfplumber → PyMuPDF → PyPDF2
  - Encrypted PDF detection with clear user message
"""

from __future__ import annotations

import csv
import importlib
from pathlib import Path
from typing import Callable


# ---------------------------------------------------------------------------
# PDF helpers
# ---------------------------------------------------------------------------

# pdfplumber table-find settings.
# Two strategies are tried in order:
#   1. "lines"      — works for PDFs with visible ruling lines (most college docs)
#   2. "text"       — works for borderless tables aligned by whitespace
_TABLE_STRATEGIES = [
    {
        "vertical_strategy":   "lines",
        "horizontal_strategy": "lines",
        "snap_tolerance":      5,
        "join_tolerance":      5,
        "edge_min_length":     10,
    },
    {
        "vertical_strategy":   "lines_strict",
        "horizontal_strategy": "lines_strict",
        "snap_tolerance":      3,
        "join_tolerance":      3,
        "edge_min_length":     10,
    },
    {
        "vertical_strategy":   "text",
        "horizontal_strategy": "lines",
        "snap_tolerance":      5,
        "join_tolerance":      5,
    },
    {
        "vertical_strategy":   "text",
        "horizontal_strategy": "text",
        "snap_tolerance":      5,
        "join_tolerance":      5,
    },
]


def _clean_cell(cell: object) -> str:
    """Normalise a table cell to a single line of text."""
    if cell is None:
        return ""
    text = str(cell).strip()
    # Collapse internal newlines that appear inside merged cells
    text = " ".join(part.strip() for part in text.splitlines() if part.strip())
    return text


def _is_empty_row(row: list[str]) -> bool:
    return not any(cell.strip() for cell in row)


def _detect_header(rows: list[list[str]]) -> int:
    """
    Return the index of the most likely header row (0 or 1).
    Heuristic: a header row has more non-empty cells than blank rows,
    and its cells tend to be shorter / title-cased.
    Returns 0 if the first row looks like a header, -1 if undetermined.
    """
    if not rows:
        return -1
    first = [c for c in rows[0] if c.strip()]
    if not first:
        return -1
    # If first row is ALL CAPS or Title Case and short, treat as header
    upper_count = sum(1 for c in first if c == c.upper() and c.strip())
    title_count = sum(1 for c in first if c.istitle())
    if upper_count >= len(first) // 2 or title_count >= len(first) // 2:
        return 0
    # Default: treat first row as header for tabular college docs
    return 0


def _format_table(raw_rows: list[list], table_index: int) -> str:
    """
    Convert raw pdfplumber table rows into a readable, retrieval-friendly string.

    Output format:
        [Table N]
        Header1 | Header2 | Header3
        --------|---------|--------
        val1    | val2    | val3
    """
    if not raw_rows:
        return ""

    # Clean every cell
    cleaned: list[list[str]] = []
    for row in raw_rows:
        row_clean = [_clean_cell(c) for c in row]
        if not _is_empty_row(row_clean):
            cleaned.append(row_clean)

    if not cleaned:
        return ""

    # Normalise column count (handle ragged rows from merged cells)
    col_count = max(len(r) for r in cleaned)
    cleaned = [r + [""] * (col_count - len(r)) for r in cleaned]

    header_idx = _detect_header(cleaned)

    lines: list[str] = [f"[Table {table_index}]"]

    for row_idx, row in enumerate(cleaned):
        line = " | ".join(row)
        lines.append(line)
        # Insert separator after header row
        if row_idx == header_idx:
            sep = " | ".join("-" * max(len(cell), 3) for cell in row)
            lines.append(sep)

    return "\n".join(lines)


def _table_bbox(table) -> tuple[float, float, float, float]:
    """Return (x0, top, x1, bottom) bounding box of a pdfplumber Table object."""
    cells = table.cells
    if not cells:
        return table.bbox
    x0  = min(c[0] for c in cells)
    top = min(c[1] for c in cells)
    x1  = max(c[2] for c in cells)
    bot = max(c[3] for c in cells)
    return x0, top, x1, bot


def _extract_page_excluding_tables(page, table_bboxes: list) -> str:
    """
    Extract prose text from a page while masking out table regions.
    This prevents table content from appearing twice (once in the table
    block and once in the raw text block).
    """
    if not table_bboxes:
        return page.extract_text(x_tolerance=3, y_tolerance=3) or ""

    # Crop the page to regions NOT covered by any table, then extract text
    remaining = page
    for bbox in table_bboxes:
        try:
            # pdfplumber: outside_bbox returns words outside a bounding box
            remaining = remaining.filter(
                lambda obj, b=bbox: not (
                    obj.get("x0", 0) >= b[0] and
                    obj.get("x1", 0) <= b[2] and
                    obj.get("top", 0) >= b[1] and
                    obj.get("bottom", 0) <= b[3]
                )
            )
        except Exception:
            pass

    return remaining.extract_text(x_tolerance=3, y_tolerance=3) or ""


def _try_extract_with_strategy(page, strategy: dict) -> list:
    """Try one table-detection strategy; return list of tables (may be empty)."""
    try:
        tables = page.find_tables(table_settings=strategy)
        return tables
    except Exception:
        return []


def _extract_tables_from_page(page) -> list:
    """
    Try table strategies in order, return the first non-empty result.
    Falls back to pdfplumber's auto-detect if all strategies fail.
    """
    for strategy in _TABLE_STRATEGIES:
        tables = _try_extract_with_strategy(page, strategy)
        if tables:
            return tables

    # Last resort: pdfplumber default auto-detect
    try:
        return page.find_tables()
    except Exception:
        return []


def _parse_pdf_pdfplumber(file_path: str) -> str:
    """Primary PDF extraction using pdfplumber with explicit table strategies."""
    import pdfplumber

    page_blocks: list[str] = []

    with pdfplumber.open(file_path) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            blocks: list[str] = []

            tables = _extract_tables_from_page(page)
            table_bboxes = [_table_bbox(t) for t in tables]

            # Prose text (table regions masked out)
            prose = _extract_page_excluding_tables(page, table_bboxes).strip()
            if prose:
                blocks.append(prose)

            # Tables
            for t_idx, table in enumerate(tables, 1):
                try:
                    raw_rows = table.extract()
                except Exception:
                    continue
                formatted = _format_table(raw_rows, t_idx)
                if formatted:
                    blocks.append(formatted)

            if blocks:
                page_blocks.append(f"[Page {page_num}]\n" + "\n\n".join(blocks))

    return "\n\n".join(page_blocks)


def _parse_pdf_pymupdf(file_path: str) -> str:
    """Secondary fallback using PyMuPDF (fitz) — better on complex/scanned PDFs."""
    fitz = importlib.import_module("fitz")  # PyMuPDF

    page_blocks: list[str] = []

    with fitz.open(file_path) as doc:
        for page_num, page in enumerate(doc, 1):
            # Extract text preserving layout (dict mode gives word-level boxes)
            blocks_raw = page.get_text("blocks", sort=True)
            texts = []
            for block in blocks_raw:
                text = block[4].strip() if len(block) > 4 else ""
                if text:
                    texts.append(text)

            # Also try to extract tables via PyMuPDF's built-in finder (v1.23+)
            try:
                table_finder = page.find_tables()
                tables = getattr(table_finder, "tables", table_finder)
                for t_idx, table in enumerate(tables, 1):
                    rows = table.extract()
                    if rows:
                        formatted = _format_table(rows, t_idx)
                        if formatted:
                            texts.append(formatted)
            except AttributeError:
                pass  # PyMuPDF < 1.23 has no find_tables

            if texts:
                page_blocks.append(f"[Page {page_num}]\n" + "\n\n".join(texts))

    return "\n\n".join(page_blocks)


def _parse_pdf_pypdf2(file_path: str) -> str:
    """Last-resort fallback using PyPDF2 (text-only, no table awareness)."""
    from PyPDF2 import PdfReader

    reader = PdfReader(file_path)
    page_blocks: list[str] = []
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        if text.strip():
            page_blocks.append(f"[Page {page_num}]\n{text.strip()}")
    return "\n\n".join(page_blocks)


def parse_pdf(file_path: str) -> str:
    """
    Extract text from a PDF with table-aware parsing.

    Strategy order:
      1. pdfplumber  — best for text-based PDFs with explicit table lines
      2. PyMuPDF     — better for complex layouts and newer PDF features
      3. PyPDF2      — plain text extraction, last resort
    """
    # Encrypted PDF check (fast, using PyPDF2's reader)
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            if not reader.decrypt(""):
                raise ValueError(
                    "PDF is password-protected. Please remove the password using "
                    "Adobe Reader or similar, then upload again."
                )
    except ValueError:
        raise
    except Exception:
        pass  # If PyPDF2 can't even open it, let the main strategies try

    # Strategy 1 — pdfplumber
    try:
        result = _parse_pdf_pdfplumber(file_path)
        if result and result.strip():
            return result
    except Exception:
        pass

    # Strategy 2 — PyMuPDF
    try:
        result = _parse_pdf_pymupdf(file_path)
        if result and result.strip():
            return result
    except ImportError:
        pass  # PyMuPDF not installed
    except Exception:
        pass

    # Strategy 3 — PyPDF2
    return _parse_pdf_pypdf2(file_path)


# ---------------------------------------------------------------------------
# Other parsers (unchanged in behaviour, minor cleanups)
# ---------------------------------------------------------------------------

def parse_docx(file_path: str) -> str:
    """Extract text from a DOCX file, preserving heading structure and tables."""
    from docx import Document

    doc = Document(file_path)
    parts: list[str] = []

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            # Paragraph
            from docx.text.paragraph import Paragraph
            para = Paragraph(block, doc)
            text = para.text.strip()
            if not text:
                continue
            style_name = para.style.name if para.style and para.style.name else ""
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading ", "").strip()
                prefix = "#" * int(level) if level.isdigit() else "##"
                parts.append(f"\n{prefix} {text}\n")
            else:
                parts.append(text)

        elif tag == "tbl":
            # Table inside the document body
            from docx.table import Table as DocxTable
            tbl = DocxTable(block, doc)
            rows: list[list[str]] = []
            for row in tbl.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            formatted = _format_table(rows, len([p for p in parts if "[Table" in p]) + 1)
            if formatted:
                parts.append(formatted)

    return "\n".join(parts)


def parse_txt(file_path: str) -> str:
    """Read a plain text file, trying common encodings before falling back."""
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            return Path(file_path).read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return Path(file_path).read_text(encoding="utf-8", errors="replace")


def parse_xlsx(file_path: str) -> str:
    """Extract text from an Excel file (all sheets), formatted as pipe-separated rows."""
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(cells)

        if rows:
            parts.append(f"\n## Sheet: {sheet_name}")
            # Reuse the same table formatter for consistent output
            formatted = _format_table(rows, 1)
            # Strip the "[Table 1]" prefix — sheet name is already the label
            body = "\n".join(formatted.splitlines()[1:]) if formatted else ""
            if body:
                parts.append(body)

    wb.close()
    return "\n".join(parts)


def parse_csv(file_path: str) -> str:
    """Extract text from a CSV file, formatted as pipe-separated rows."""
    rows: list[list[str]] = []
    with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as f:
        for row in csv.reader(f):
            if any(cell.strip() for cell in row):
                rows.append(row)
    return _format_table(rows, 1) if rows else ""


def parse_pptx(file_path: str) -> str:
    """Extract text from a PowerPoint file, slide by slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is not None:
                for para in text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))

    return "\n\n".join(parts)


def parse_html(file_path: str) -> str:
    """Extract text from an HTML file, stripping navigation and boilerplate."""
    from bs4 import BeautifulSoup

    content = Path(file_path).read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(content, "lxml")

    for tag in soup(["script", "style", "nav", "footer", "header", "noscript"]):
        tag.decompose()

    lines = [line.strip() for line in soup.get_text(separator="\n").splitlines() if line.strip()]
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Registry & public API
# ---------------------------------------------------------------------------

PARSERS: dict[str, Callable[[str], str]] = {
    ".pdf":  parse_pdf,
    ".docx": parse_docx,
    ".txt":  parse_txt,
    ".xlsx": parse_xlsx,
    ".csv":  parse_csv,
    ".pptx": parse_pptx,
    ".html": parse_html,
    ".htm":  parse_html,
}


def parse_document(file_path: str) -> tuple[str | None, str | None]:
    """
    Parse a document and return its extracted text.

    Returns:
        (text, None)           on success
        (None, error_message)  on failure
    """
    path = Path(file_path)
    ext  = path.suffix.lower()
    parser = PARSERS.get(ext)

    if not parser:
        supported = ", ".join(sorted(PARSERS))
        return None, f"Unsupported file type '{ext}'. Supported: {supported}"

    try:
        text = parser(file_path)
        if not text or not text.strip():
            return None, "Document appears to be empty or contains no extractable text."
        return text.strip(), None

    except ValueError as exc:
        # e.g. encrypted PDF — surface the message directly to the user
        return None, str(exc)

    except Exception as exc:
        return None, f"Failed to parse {path.name}: {exc}"